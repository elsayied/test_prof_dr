# pptx_utils.py
from pptx import Presentation
from pptx.util import Inches
from typing import List, Dict
import re
import os
from google import genai

# --- HARDCODED API KEY (SECURITY WARNING: DO NOT USE IN PUBLIC REPOSITORIES) ---
# ⚠️ REPLACE THE PLACEHOLDER BELOW WITH YOUR ACTUAL GEMINI API KEY ⚠️
HARDCODED_GEMINI_API_KEY = "AIzaSyA4YsTnbNjl2gKn20EqPa-9nom9yymEwd0"
# ---

# --- Docling Setup ---
# Check for docling availability and set up extraction function or mock
try:
    from docling.document_converter import DocumentConverter
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import (
        PdfPipelineOptions,
        TesseractCliOcrOptions,
    )
    from docling.document_converter import PdfFormatOption
    DOCLING_AVAILABLE = True
except ImportError:
    DOCLING_AVAILABLE = False
    
    # Mock function definition for unavailable docling
    def extract_content_with_docling(file_path: str, page_range: str = None) -> str:
        """Mocks docling extraction when dependencies are missing."""
        return f"# Mock Extracted Content from {os.path.basename(file_path)}\n\nThis is mock content because docling is not installed or available.\n\n- Key Document Section 1\n  - Detail A\n- Key Document Section 2\n"

if DOCLING_AVAILABLE:
    def extract_content_with_docling(file_path: str, page_range: str = None) -> str:
        """Extracts content from various file types using the docling library."""
        try:
            ocr_options = TesseractCliOcrOptions(lang=["eng"])
            pipeline_options = PdfPipelineOptions(
                do_ocr=True, do_table_structure=True, ocr_options=ocr_options
            )

            doc_converter = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(
                        pipeline_options=pipeline_options,
                    )
                }
            )
            doc = doc_converter.convert(file_path).document
            
            if page_range:
                pages = []
                for part in page_range.split(','):
                    part = part.strip()
                    if '-' in part:
                        start, end = map(int, part.split('-'))
                        pages.extend(range(start - 1, end))
                    elif part.isdigit():
                        pages.append(int(part) - 1)
                
                filtered_content = []
                for i, page in enumerate(doc.pages):
                    if i in pages:
                        filtered_content.append(page.export_to_markdown())
                return "\n".join(filtered_content)
            else:
                return doc.export_to_markdown()
        
        except Exception as e:
            raise Exception(f"Docling extraction error: {e}")


# --- Gemini LLM Integration ---

GEMINI_LECTURE_PROMPT = '''
You are an expert instructional designer and content synthesiser. Your task is to take the provided raw document text, synthesize the key information, and structure it into a series of slides using Markdown format.

Your output MUST be a single Markdown string.

### Output Format Rules:
1.  **Slide Separation:** Use three hyphens (`---`) on a new line to separate each slide.
2.  **Titles:** Use `# ` for the slide title.
3.  **Layout Hint:** Start the content of the first slide with `layout: title_slide` on a new line after the title. For subsequent slides, use `layout: title_content` or `layout: comparison` if appropriate, followed by the slide title.
4.  **Content:** Use standard Markdown list items (`-`, `  -`, etc.) for bullet points. Maintain proper indentation (two spaces per level) for sub-points.
5.  **Tables:** Use standard GitHub-flavored Markdown for tables, including a header and separator line (e.g., `| Col1 | Col2 |` followed by `|---|---|`).
6.  **Comparison Slides:** For `layout: comparison` slides, separate the content for the left and right columns with `|||` on a new line.

Here is the raw text you need to convert:
'''

def generate_structured_markdown(raw_text: str, model: str = 'gemini-2.5-flash') -> str:
    """
    Calls the Gemini API to convert raw text into structured slide Markdown,
    using a hardcoded key if available.
    """
    api_key = HARDCODED_GEMINI_API_KEY
    if api_key == "YOUR_HARDCODED_GEMINI_API_KEY_HERE" or not api_key:
        # Fallback to local mock if key is missing or is the placeholder
        return generate_structured_markdown_mock(raw_text, is_hardcoded_missing=True)

    # Use the hardcoded key directly
    genai.configure(api_key=api_key)
    client = genai.Client()
    
    try:
        response = client.models.generate_content(
            model=model,
            contents=GEMINI_LECTURE_PROMPT + raw_text
        )
        return response.text
    except APIError as e:
        # If API call fails (e.g., key invalid, quota exceeded)
        return f"# Gemini API Error\n- Failed to generate content: {e}\n- Please check your API Key."
    except Exception as e:
        return f"# General Error\n- An unexpected error occurred: {e}"

def generate_structured_markdown_mock(raw_text: str, is_hardcoded_missing: bool = False) -> str:
    """
    Mock function used when no API key is available or hardcoded key is missing.
    """
    warning_text = "API Key MISSING or is the placeholder." if is_hardcoded_missing else "Docling is mocked."
    title_line = raw_text.split('\n')[0].strip() or "Untitled Document"
    
    return f"""
layout: title_slide
# LLM MOCK: {title_line}
- Status: Fallback Mock Content ({warning_text})
- Generated from {len(raw_text)} characters

---

layout: title_content
# Key Concepts (MOCKED)
- Concept 1 synthesized from the first paragraph.
  - Detail A supporting the idea.
- Concept 2 from the main body.

---

layout: comparison
# Comparison Table (MOCKED)
Feature A ||| Feature B
- Value 1 ||| - Value X
- Value 2 ||| - Value Y

---

layout: title_content
# Summary
This content is a placeholder. Insert your key in pptx_utils.py for real generation.
"""

# --- Markdown Parsing and PPTX Creation Functions (Unchanged) ---

def create_table_from_markdown(text: str) -> List[List[str]]:
    """Convert Markdown table to table data."""
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    if not lines: return []
    lines = [l for l in lines if not re.match(r'^\s*\|?.*--.*\|?\s*$', l)]
    
    table_data = []
    for row_str in lines:
        row_str = row_str.strip()
        if row_str.startswith('|'): row_str = row_str[1:]
        if row_str.endswith('|'): row_str = row_str[:-1]
        cells = [cell.strip() for cell in row_str.split('|')]
        table_data.append(cells)
    return table_data

def add_bullet_points_from_markdown(text_frame, points: str):
    """Add bullet points to a text frame from Markdown list."""
    if not text_frame.text.strip(): text_frame.text = ""

    def get_level_and_text(line: str) -> tuple[int, str]:
        stripped_line = line.lstrip()
        text = stripped_line
        if text.startswith(('-', '*', '+')) and text[1:2] in (' ', ''):
            text = text[1:].lstrip()
        indent = len(line) - len(line.lstrip())
        level = indent // 2
        return level, text

    lines = [line for line in points.split('\n') if line.strip()]
    if not lines: return

    for line in lines:
        level, text = get_level_and_text(line)
        p = text_frame.add_paragraph()
        p.text = text
        p.level = min(level, 8)

def parse_markdown_to_slides(content: str) -> List[Dict]:
    """Parses markdown-like text into a list of slide definitions."""
    slides = []
    slide_contents = re.split(r'\n---\n', content)

    for slide_content in slide_contents:
        if not slide_content.strip(): continue

        lines = slide_content.strip().split('\n')
        slide_data = {'layout': 'title_content', 'title': None, 'blocks': []}

        if lines and lines[0].strip().lower().startswith('layout:'):
            slide_data['layout'] = lines[0].split(':', 1)[1].strip()
            lines.pop(0)
        
        if lines and lines[0].strip().startswith('# '):
            slide_data['title'] = lines[0][2:].strip()
            lines.pop(0)

        line_idx = 0
        while line_idx < len(lines):
            line = lines[line_idx]
            if not line.strip(): line_idx += 1; continue

            if line.lstrip().startswith(('-', '*', '+')):
                bullet_lines = []
                start_indent = len(line) - len(line.lstrip())
                while line_idx < len(lines):
                    current_line = lines[line_idx]
                    current_indent = len(current_line) - len(current_line.lstrip())
                    
                    if current_line.strip() and current_indent < start_indent and not current_line.lstrip().startswith(('-', '*', '+')):
                        break
                    
                    if not current_line.strip() and line_idx + 1 < len(lines) and lines[line_idx+1].strip() and (len(lines[line_idx+1]) - len(lines[line_idx+1].lstrip()) < start_indent):
                        break
                        
                    if current_line.strip() or current_indent >= start_indent:
                        bullet_lines.append(current_line)
                    line_idx += 1
                
                slide_data['blocks'].append({'type': 'bullet', 'content': '\n'.join(bullet_lines)})
                continue

            is_table = False
            if line.strip().startswith('|'):
                if (line_idx + 1 < len(lines)) and re.match(r'^\s*\|?.*--.*\|?\s*$', lines[line_idx+1]):
                    is_table = True
            
            if is_table:
                table_lines = []
                table_lines.append(lines[line_idx])
                line_idx += 1
                if line_idx < len(lines):
                    table_lines.append(lines[line_idx])
                    line_idx += 1

                while line_idx < len(lines) and lines[line_idx].strip().startswith('|'):
                    table_lines.append(lines[line_idx])
                    line_idx += 1
                    
                slide_data['blocks'].append({'type': 'table', 'content': '\n'.join(table_lines)})
                continue

            text_lines = []
            while line_idx < len(lines):
                current_line = lines[line_idx]
                if not current_line.strip():
                    if line_idx + 1 < len(lines) and not lines[line_idx+1].lstrip().startswith(('-', '*', '+')) and not lines[line_idx+1].strip().startswith('|'):
                        text_lines.append(current_line)
                        line_idx += 1
                        continue
                    else:
                        break
                
                if current_line.lstrip().startswith(('-', '*', '+')): break
                if current_line.strip().startswith('|') and (line_idx + 1 < len(lines)) and re.match(r'^\s*\|?.*--.*\|?\s*$', lines[line_idx+1]): break
                    
                text_lines.append(current_line)
                line_idx += 1
            
            if text_lines:
                slide_data['blocks'].append({'type': 'text', 'content': '\n'.join(text_lines)})
            
            if not text_lines and not is_table and not (line.lstrip().startswith(('-', '*', '+'))):
                line_idx += 1
            
        slides.append(slide_data)
    return slides

def create_presentation_from_markdown(content: str, output_path: str = 'output.pptx') -> str:
    """
    Create a PowerPoint presentation from Markdown-structured text.
    (Contains the logic to map Markdown to PPTX slides and elements)
    """
    prs = Presentation()
    layout_map = {
        'title_slide': prs.slide_layouts[0], 'title_content': prs.slide_layouts[1],
        'section_header': prs.slide_layouts[2], 'two_content': prs.slide_layouts[3],
        'comparison': prs.slide_layouts[4], 'title_only': prs.slide_layouts[5],
        'blank': prs.slide_layouts[6],
    }

    slides_data = parse_markdown_to_slides(content)

    for slide_data in slides_data:
        layout_name = slide_data['layout']
        slide_layout = layout_map.get(layout_name, layout_map['title_content'])
        current_slide = prs.slides.add_slide(slide_layout)

        if slide_data['title'] and current_slide.shapes.title:
            current_slide.shapes.title.text = slide_data['title']

        if layout_name in ['comparison', 'two_content']:
            content_placeholders = [p for p in current_slide.placeholders if p.placeholder_format.idx > 0 and p.has_text_frame]
            left_text_block = next((block for block in slide_data['blocks'] if block['type'] == 'text'), None)
            
            if left_text_block and len(content_placeholders) >= 2:
                left_ph, right_ph = content_placeholders[0], content_placeholders[1]
                text_block_content = left_text_block['content']
                
                left_text, right_text = text_block_content.split('|||', 1) if '|||' in text_block_content else (text_block_content, "")
                
                left_ph.text_frame.clear()
                right_ph.text_frame.clear()
                left_ph.text_frame.paragraphs[0].text = left_text.strip()
                right_ph.text_frame.paragraphs[0].text = right_text.strip()
            
        else:
            body_shape = next((shape for shape in current_slide.placeholders if shape.placeholder_format.idx != 0 and shape.has_text_frame), None)
            
            if body_shape:
                tf = body_shape.text_frame
                tf.clear()
                
                if len(tf.paragraphs) > 0 and not tf.paragraphs[0].text.strip():
                     tf.paragraphs[0].clear()

                for block in slide_data['blocks']:
                    if block['type'] == 'text':
                        p = tf.add_paragraph()
                        p.text = block['content']
                    elif block['type'] == 'bullet':
                        add_bullet_points_from_markdown(tf, block['content'])
                    elif block['type'] == 'table':
                        table_data = create_table_from_markdown(block['content'])
                        if not table_data: continue

                        rows, cols = len(table_data), len(table_data[0])
                        table_width = Inches(8.0)
                        table_height = Inches(0.4 * (rows + 1))
                        left = (prs.slide_width - table_width) / 2
                        top = Inches(2.0)
                        
                        try:
                            table_shape = current_slide.shapes.add_table(rows, cols, left, top, table_width, table_height)
                            table = table_shape.table
                            for r_idx, row_data in enumerate(table_data):
                                for c_idx, cell_text in enumerate(row_data):
                                    if c_idx < cols:
                                        table.cell(r_idx, c_idx).text = cell_text
                        except ValueError:
                            p = tf.add_paragraph()
                            p.text = "[ERROR: Could not create table due to inconsistent data. Check Markdown format.]"
                            
    prs.save(output_path)
    return output_path
